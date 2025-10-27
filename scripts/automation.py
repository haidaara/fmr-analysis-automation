import os
import yaml
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from loading_setup import loading_setup
from fit_spectrum import fit_spectrum

import logging

# Utilities moved to separate module for clarity
from utils import (
    discover_raw_files,
    simple_parse_filename,
    archive_and_prepare_sample_dir,
    collect_fit_results,
    collect_plot_files,
    collect_fit_plot_files,
    filtered_pptx_formats,
    validate_image_or_warn,
)

# ========== CONFIGURATION LOADER ==========

def load_config(config_path="config.yaml", cli_args=None):
    """
    Load configuration from YAML file and merge CLI arguments.
    CLI arguments override YAML keys.
    """
    with open(config_path, "r") as f:
        config = yaml.safe_load(f)
    if cli_args:
        for k, v in cli_args.items():
            if v not in [None, "", [], {}]:
                config[k] = v
    # Set defaults if not present
    config.setdefault("data_folder", "data")
    config.setdefault("results_root", "results_roots")
    config.setdefault("samples", [])
    config.setdefault("plot_format", "png")
    config.setdefault("aggregation", "median")
    config.setdefault("model", "double")
    return config

# ========== REPORT ARTIFACTS ==========

def create_sample_analysis_xlsx(sample, results_root):
    rows = collect_fit_results(sample, results_root)
    if not rows:
        logging.warning(f"No fit results for sample {sample}")
        print("\n" + "-"*60)
        print(f"[ANALYSIS] No fit results for sample: {sample}")
        print("-"*60 + "\n")
        return
    meta_cols = ["sample", "f", "db", "index", "result_file"]
    data_cols = [c for c in rows[0] if c not in meta_cols]
    columns = meta_cols + data_cols
    df = pd.DataFrame(rows)
    df = df[columns]
    out_path = os.path.join(results_root, sample, f"{sample}_analysis.xlsx")
    df.to_excel(out_path, index=False)
    logging.info(f"[{sample}] Aggregated results saved to {out_path}")
    print("\n" + "-"*60)
    print(f"[ANALYSIS] Aggregated results saved to: {out_path}")
    print("-"*60 + "\n")

def create_sample_pptx(sample, results_root, plot_formats=["png"]):
    # PPTX supports raster-only; filter formats here
    formats = filtered_pptx_formats(plot_formats)
    plot_files = []
    for fmt in formats:
        plot_files.extend(collect_plot_files(sample, results_root, [fmt]))
    plot_files = sorted(plot_files)

    if not plot_files:
        logging.warning(f"No plot files found for {sample}")
        print("\n" + "-"*60)
        print(f"[PPTX] No data-only plot files found for sample: {sample}")
        print("-"*60 + "\n")
        return

    print("\n" + "-"*60)
    print(f"[PPTX] Creating plots PPTX for sample: {sample}")
    print("-"*60)

    prs = Presentation()
    for plot_path in plot_files:
        if not validate_image_or_warn(plot_path):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        fname = os.path.basename(plot_path)
        meta = simple_parse_filename(fname)
        if meta is None:
            continue
        db_val = meta.get('db', '')
        title_text = f"{meta.get('sample','')} f={meta.get('f','')}GHz {db_val} dB"
        left = Inches(0.5)
        top = Inches(0.5)
        slide.shapes.add_picture(plot_path, left, top, width=Inches(8), height=Inches(5))
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.3))
        tf = txBox.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        print(f"\t[ADD] {fname}")

    pptx_path = os.path.join(results_root, sample, f"{sample}_plots.pptx")
    prs.save(pptx_path)
    logging.info(f"[{sample}] PowerPoint with plots saved to {pptx_path}")
    print(f"[PPTX] Saved: {pptx_path}")
    print("-"*60 + "\n")

def create_sample_fit_pptx(sample, results_root, plot_formats=["png"]):
    formats = filtered_pptx_formats(plot_formats)
    fit_plot_files = []
    for fmt in formats:
        fit_plot_files.extend(collect_fit_plot_files(sample, results_root, [fmt]))
    fit_plot_files = sorted(fit_plot_files)

    if not fit_plot_files:
        logging.warning(f"No fit plot files found for {sample}")
        print("\n" + "-"*60)
        print(f"[PPTX] No fit plot files found for sample: {sample}")
        print("-"*60 + "\n")
        return

    print("\n" + "-"*60)
    print(f"[PPTX] Creating fit PPTX for sample: {sample}")
    print("-"*60)

    prs = Presentation()
    for plot_path in fit_plot_files:
        if not validate_image_or_warn(plot_path):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        fname = os.path.basename(plot_path)
        meta = simple_parse_filename(fname)
        if meta is None:
            continue
        db_val = meta.get('db', '')
        title_text = f"{meta.get('sample','')} f={meta.get('f','')}GHz {db_val} dB FIT"
        left = Inches(0.5)
        top = Inches(0.5)
        slide.shapes.add_picture(plot_path, left, top, width=Inches(8), height=Inches(5))
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.3))
        tf = txBox.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        print(f"\t[ADD] {fname}")

    pptx_path = os.path.join(results_root, sample, f"{sample}_fit.pptx")
    prs.save(pptx_path)
    logging.info(f"[{sample}] PowerPoint with fit plots saved to {pptx_path}")
    print(f"[PPTX] Saved: {pptx_path}")
    print("-"*60 + "\n")

# ========== MAIN PIPELINE ==========

def main():
    import argparse
    parser = argparse.ArgumentParser(description="Run FMR automation pipeline.")
    parser.add_argument("--config", type=str, default="config.yaml", help="Path to YAML configuration file.")
    parser.add_argument("--log_level", type=str, default="", help="Logging level: INFO, DEBUG, WARNING, ERROR.")
    args = parser.parse_args()
    cli_args = vars(args)
    config = load_config(cli_args.pop("config"), cli_args)

    # File-only logging (console uses clean prints)
    log_level = getattr(logging, str(config.get("log_level", "INFO")).upper(), logging.INFO)
    os.makedirs(config["results_root"], exist_ok=True)
    logging.basicConfig(
        level=log_level,
        format="%(asctime)s [%(levelname)s] %(message)s",
        handlers=[
            logging.FileHandler(os.path.join(config["results_root"], "automation.log"))
        ]
    )

    data_folder = config["data_folder"]
    results_root = config["results_root"]
    sample_list = config["samples"] if config["samples"] else None
    plot_formats = config.get("plot", {}).get("plot_formats", [config.get("plot_format", "png")]) \
        if "plot" in config else [config.get("plot_format", "png")]

    discovered = discover_raw_files(data_folder, sample_list)
    if not discovered:
        logging.warning("No raw files discovered.")
        print("\n" + "="*60)
        print("[INFO] No raw files discovered.")
        print("="*60 + "\n")
        return

    samples = []
    for d in discovered:
        if d["sample"] not in samples:
            samples.append(d["sample"])

    for sample in samples:
        print("\n" + "="*70)
        print(f"Processing sample: {sample}")
        print("-"*70)

        logging.info(f"=== Processing sample: {sample} ===")
        sample_discovered = [d for d in discovered if d["sample"] == sample]

        # Archive existing sample folder and prepare a fresh output folder
        archive_and_prepare_sample_dir(results_root, sample)

        sample_discovered = sorted(
            sample_discovered,
            key=lambda d: (
                d['f_float'],
                d['db'] if d['db'] is not None else 0,
                d.get('index', 0) or 0
            )
        )

        seen_paths = set()
        for meta in sample_discovered:
            if meta['full_path'] in seen_paths:
                continue
            seen_paths.add(meta['full_path'])

            logging.info(f"-> Processing: {meta['filename']}")
            print(f"\t-> Processing: {meta['filename']}")

            processed_path = loading_setup(
                os.path.dirname(meta['full_path']),
                results_root,
                meta['sample'],
                meta['f'],
                meta['db']
            )

            fit_config = config.copy()
            fit_config["input_file"] = processed_path
            fit_config["sample"] = meta["sample"]
            fit_config["results_root"] = results_root
            fit_spectrum(cli_args=fit_config)

        create_sample_analysis_xlsx(sample, results_root)
        create_sample_pptx(sample, results_root, plot_formats=plot_formats)
        create_sample_fit_pptx(sample, results_root, plot_formats=plot_formats)

        print("-"*70)
        print(f"DONE sample: {sample}")
        print("="*70 + "\n")

if __name__ == "__main__":
    main()